import functools
import itertools
import operator
import os
import re
from statistics import mean
from typing import cast

# import fitz as pymupdf
import pymupdf
from dotenv import load_dotenv
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation
from pptx.shapes.autoshape import Shape

load_dotenv(override=True)

DATA_DIR = os.getenv("BOOKS_DIR", "")


# ----------------------------------------------------------------------------
# Flatten chunks from dictionary to list
# ----------------------------------------------------------------------------

def flatten_chunks(documents_dict: dict[str, list[str]]) -> list[str]:

    values_list = list(documents_dict.values())
    flattened_list = list(itertools.chain(*values_list))

    return flattened_list


# ----------------------------------------------------------------------------
# Flatten chunks from dictionary to list
# ----------------------------------------------------------------------------

DEFAULT_CHUNK_SIZE = 1100
DEFAULT_CHUNK_OVERLAP = 200


# -----------------------------------------------------------------------------
# Read files from directory and subdirectories
# -----------------------------------------------------------------------------

def read_files_directory(base_path: str = DATA_DIR,
                         chunk_size: int = DEFAULT_CHUNK_SIZE,
                         chunk_overlap: int = DEFAULT_CHUNK_OVERLAP
                         ) -> tuple[list[str], int, list[str]]:
    """
    This function takes a directory path as input and reads every file
    in the directory and its subdirectories using the ReadFile function.

    Args:
    directory_path (str): The path to the directory containing the files.

    Returns:
    tuple: The flattened chunks, number of files found, and files that could
    not be read.
    """

    loader = DocumentIngest(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    print(f'Document loader: {loader}')

    files_content = {}
    files_count = 0
    problem_files = []

    # Walk through the directory tree
    for root, _, files in os.walk(base_path):
        print(f"Reading files in {root}")
        # print(f"Files: {files}")
        for filename in files:
            files_count += 1
            # print(f"Reading {filename}")
            filepath = os.path.join(root, filename)

            try:
                # Read the file using the ReadFile function
                content = loader.load_and_split(filepath)
                files_content[filepath] = content
                # print(f"Reading and storing {filepath}")
            except Exception as err:  # noqa: BLE001
                print(f"Error reading file {filepath}: {err}")
                problem_files.append(filepath)

    filtered_dict = {k: v for k, v in files_content.items() if v is not None}

    values_list = list(filtered_dict.values())
    flattened_list = list(itertools.chain(*values_list))

    return flattened_list, files_count, problem_files


# ----------------------------------------------------------------------------
# Document Ingest class
# ----------------------------------------------------------------------------

class DocumentIngest:

    def __init__(self,
                 chunk_size: int = 1100,
                 chunk_overlap: int = 200):
        self.text_splitter = \
            self.create_text_splitter(chunk_size=chunk_size,
                                      chunk_overlap=chunk_overlap)

    def create_text_splitter(self, chunk_size: int, chunk_overlap: int
                             ) -> RecursiveCharacterTextSplitter:
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            is_separator_regex=False,
            separators=["\n\n"]
        )

        return text_splitter

    @staticmethod
    def get_mean_line_length(text: str) -> float:
        lines = text.split('\n')
        word_count = [len(i.split(' ')) for i in lines]
        return mean(word_count)

    @staticmethod
    def read_pdf_as_text(document_path: str) -> str:

        doc = pymupdf.open(document_path)
        n_pages = doc.page_count

        # Extract all blocks from all pages
        list_pages = [doc[i].get_text("blocks") for i in range(n_pages)]

        # Filter by block
        list_pages_clean_blocks = [
            DocumentIngest.filter_blocks(
                cast(
                    list[tuple[float, float, float, float, str, int, int]],
                    blocks,
                )
            )
            for blocks in list_pages
        ]

        # Join all blocks per page
        list_pages_clean_str = [''.join(i) for i in list_pages_clean_blocks]

        # Filter by mean line length
        list_pages_clean_str_filter = [
            page for page in list_pages_clean_str
            if DocumentIngest.get_mean_line_length(page)]

        # Filter by page content
        list_pages_superclean = DocumentIngest.filter_pages(
            list_pages_clean_str_filter)

        # Join all pages
        text = '\n'.join(list_pages_superclean)

        return text

    @staticmethod
    def read_pptx_as_text(document_path: str) -> str:
        presentation = Presentation(document_path)
        extracted_texts: list[str] = []

        for slide in presentation.slides:

            # Collect the texts from the components (shapes) of the slide
            text = []
            for shape in slide.shapes:
                if isinstance(shape, Shape):
                    shape_text = shape.text_frame.text.strip()
                    if shape_text != '':
                        text.append(shape_text)

            # Remove trailing slide number from texts
            text = text[:-1] if len(text) > 0 else text

            # Append the extracted text to the list of extracted texts
            if text:
                extracted_texts.append(' '.join(text))
            else:
                # Insert None placeholders so we know slide number
                # (i.e. index) and number of empty slides
                extracted_texts.append('')

        # Return the extracted text from all slides as a single string
        extracted_text = ' '.join(extracted_texts)
        extracted_text = extracted_text.replace('\n', ' ')

        return extracted_text

    @staticmethod
    def filter_pages(list_pages: list[str]) -> list[str]:

        list_pages_clean = [page for page in list_pages if len(page) > 0]
        list_pages_clean = [page for page in list_pages_clean if not bool(
            re.search("Intentionally Blank", page))]

        return list_pages_clean

    @staticmethod
    def filter_blocks(
        list_blocks: list[
            tuple[float, float, float, float, str, int, int]
        ],
        y_diff: int = 400,
        y_min: int = 55,
        y_max: int = 735,
        min_len: int = 5,
    ) -> list[str]:

        list_blocks_clean = [i for i in list_blocks if i[3]-i[1] < y_diff]
        list_blocks_clean = [i for i in list_blocks_clean if i[3] > y_min]
        list_blocks_clean = [i for i in list_blocks_clean if i[3] < y_max]
        list_blocks_clean = [i for i in list_blocks_clean
                             if len(i[4]) > min_len]

        list_text = [i[4] for i in list_blocks_clean]
        list_text_clean = [
            i for i in list_text if not bool(re.search('^Figure', i))]
        list_text_clean = [
            i for i in list_text_clean if not bool(re.search('^Table', i))]

        if len(list_text_clean) > 1:
            return list_text_clean
        else:
            return []

    def load_and_split(self, document_path: str) -> list[str]:
        # regex = (\d{7})_{0,1}\d{0,2}\.(pdf|pptx)

        try:
            extension = os.path.splitext(document_path)[1].lower()
            if extension == '.pdf':
                text = DocumentIngest.read_pdf_as_text(document_path)
                # print (f"{document_path} document has been split and loaded")
            elif extension == '.pptx':
                text = DocumentIngest.read_pptx_as_text(document_path)
            else:
                # print (f"{document_path} format not taken into account,
                #  document is not parsed")
                return []

            text_chunks = self.text_splitter.split_text(text)

            return text_chunks

        except Exception as err:  # noqa: BLE001
            print(f"Error reading {document_path}, document not parsed\n{err}")
            return []


# ----------------------------------------------------------------------------
# Group texts into blocks
# ----------------------------------------------------------------------------


def group_texts(
    examples: dict[str, list[list[int]]],
    block_size: int = 256,
) -> dict[str, list[list[int]]]:
    """
    This column is not present in ipynb but causes an error in script.
    Removing it seems to fix issue.
    TODO: Understand bug and fix it properly.
    """
    if 'overflow_to_sample_mapping' in examples:
        examples.pop('overflow_to_sample_mapping')

    # Concatenate all texts.
    concatenated_examples = \
        {k: functools.reduce(operator.iadd, examples[k], [])
         for k in examples}
    total_length = \
        len(concatenated_examples[next(iter(examples.keys()))])

    # We drop the small remainder, we could add padding if the model supported
    # it instead of this drop, you customize this part to your needs.
    total_length = (total_length // block_size) * block_size

    # Split by chunks of max_len.
    result = {
        k: [t[i: i + block_size] for i in range(0, total_length, block_size)]
        for k, t in concatenated_examples.items()
    }

    result["labels"] = result["input_ids"].copy()

    return result

# ----------------------------------------------------------------------------
# End of File
# ----------------------------------------------------------------------------
